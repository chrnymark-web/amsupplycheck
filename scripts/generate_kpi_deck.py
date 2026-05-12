"""Generate SupplyCheck Weekly KPI Deck (.pptx).

Output: SupplyCheck_Weekly_KPI_Deck.pptx in project root.
Branding follows skills/slide-deck/SKILL.md (dark header + green accent + Georgia/Arial).
"""

from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
LOGO_DARK = str(ROOT / "src" / "assets" / "amsupplycheck-logo-new.png")
LOGO_WHITE = str(ROOT / "src" / "assets" / "amsupplycheck-logo-white.png")
OUTPUT = ROOT / "SupplyCheck_Weekly_KPI_Deck.pptx"

GREEN = RGBColor(0x6B, 0x8F, 0x3C)
GREEN_DARK = RGBColor(0x4F, 0x6C, 0x2A)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_LIGHT = RGBColor(0xF2, 0xF2, 0xF2)
GRAY_MED = RGBColor(0xDD, 0xDD, 0xDD)
GRAY_TEXT = RGBColor(0x55, 0x55, 0x55)
DARK_TEXT = RGBColor(0x33, 0x33, 0x33)
RED_MUTED = RGBColor(0xCC, 0x44, 0x44)
AMBER = RGBColor(0xE0, 0xA0, 0x30)

FONT_HEAD = "Georgia"
FONT_BODY = "Arial"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height


def no_border(shape):
    shape.line.fill.background()


def fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color, shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, left, top, width, height)
    fill(s, color)
    no_border(s)
    return s


def add_text(
    slide,
    left,
    top,
    width,
    height,
    text,
    size=14,
    color=DARK_TEXT,
    bold=False,
    font=FONT_BODY,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    return tf


def add_runs(
    slide,
    left,
    top,
    width,
    height,
    runs,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    line_spacing=None,
):
    """runs: list of (text, size, color, bold, font) tuples. '\n' in text starts a new paragraph."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    first = True
    for text, size, color, bold, font in runs:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = font
    return tf


def add_header(slide, title_text, slide_num):
    # Dark header bar
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.4), DARK)
    # Title
    add_text(
        slide,
        Inches(0.5),
        Inches(0.3),
        Inches(10),
        Inches(0.9),
        title_text,
        size=32,
        color=WHITE,
        bold=True,
        font=FONT_HEAD,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    # Green accent strip
    add_rect(slide, Inches(12.333), 0, Inches(1), Inches(1.4), GREEN)
    # Slide number
    add_text(
        slide,
        Inches(12.333),
        Inches(0.25),
        Inches(1),
        Inches(0.9),
        str(slide_num),
        size=32,
        color=WHITE,
        bold=True,
        font=FONT_HEAD,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    # Logo top-right (inside dark header, before accent strip)
    slide.shapes.add_picture(
        LOGO_WHITE, Inches(11.0), Inches(0.45), height=Inches(0.5)
    )


def add_bg(slide):
    add_rect(slide, 0, Inches(1.4), SLIDE_W, Inches(6.1), GRAY_LIGHT)


def blank_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])


# =====================================================================
# Slide 1 — Title
# =====================================================================
def slide_title():
    s = blank_slide()
    # Light left panel
    add_rect(s, 0, 0, Inches(9.333), SLIDE_H, GRAY_LIGHT)
    # Dark right panel
    add_rect(s, Inches(9.333), 0, Inches(4), SLIDE_H, DARK)
    # Green accent strip down the middle seam
    add_rect(s, Inches(9.233), 0, Inches(0.1), SLIDE_H, GREEN)

    # Logo (dark version) on light bg
    s.shapes.add_picture(LOGO_DARK, Inches(0.7), Inches(0.7), height=Inches(0.9))

    # Title
    add_runs(
        s,
        Inches(0.7),
        Inches(2.4),
        Inches(8.2),
        Inches(2.0),
        [
            ("Vejen til 30k/md\n", 52, DARK_TEXT, True, FONT_HEAD),
            ("Ugentlige KPI'er for SupplyCheck", 26, GRAY_TEXT, False, FONT_BODY),
        ],
        line_spacing=1.1,
    )

    # Subtitle / north star
    add_runs(
        s,
        Inches(0.7),
        Inches(5.0),
        Inches(8.2),
        Inches(1.5),
        [
            ("10.000 mdr. besøgende  →  30.000 DKK/md i løn\n", 20, GREEN_DARK, True, FONT_BODY),
            ("Personligt accountability-deck · review hver mandag morgen", 14, GRAY_TEXT, False, FONT_BODY),
        ],
        line_spacing=1.4,
    )

    # Date in dark panel
    today = date.today().strftime("%d. %B %Y")
    add_text(
        s,
        Inches(9.5),
        Inches(6.7),
        Inches(3.5),
        Inches(0.5),
        today,
        size=12,
        color=WHITE,
        font=FONT_BODY,
        align=PP_ALIGN.CENTER,
    )

    # Big checkmark on dark panel (using a green circle + white text)
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.4), Inches(2.7), Inches(2), Inches(2))
    fill(circle, GREEN)
    no_border(circle)
    add_text(
        s,
        Inches(10.4),
        Inches(2.7),
        Inches(2),
        Inches(2),
        "✓",
        size=88,
        color=WHITE,
        bold=True,
        font=FONT_HEAD,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )


# =====================================================================
# Slide 2 — North Star
# =====================================================================
def slide_north_star():
    s = blank_slide()
    add_header(s, "North Star", 2)
    add_bg(s)

    # Big equation centered
    add_text(
        s,
        Inches(0.5),
        Inches(2.0),
        Inches(12.333),
        Inches(0.6),
        "NORTH STAR",
        size=14,
        color=GRAY_TEXT,
        bold=True,
        font=FONT_BODY,
        align=PP_ALIGN.CENTER,
    )

    # Three boxes: 10k visitors  →  conversion  →  30k DKK
    box_w = Inches(3.6)
    box_h = Inches(2.4)
    gap = Inches(0.4)
    total_w = box_w * 3 + gap * 2
    start_x = (SLIDE_W - total_w) / 2
    y = Inches(2.8)

    # Box 1
    b1 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, start_x, y, box_w, box_h)
    fill(b1, WHITE)
    no_border(b1)
    add_runs(
        s,
        start_x,
        y,
        box_w,
        box_h,
        [
            ("10.000\n", 56, GREEN_DARK, True, FONT_HEAD),
            ("månedlige besøgende\n", 14, GRAY_TEXT, False, FONT_BODY),
            ("buyer-trafik på platformen", 11, GRAY_TEXT, False, FONT_BODY),
        ],
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
        line_spacing=1.2,
    )

    # Arrow 1
    arrow_y = y + box_h / 2 - Inches(0.25)
    a1 = s.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        start_x + box_w,
        arrow_y,
        gap,
        Inches(0.5),
    )
    fill(a1, GREEN)
    no_border(a1)

    # Box 2
    x2 = start_x + box_w + gap
    b2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x2, y, box_w, box_h)
    fill(b2, DARK)
    no_border(b2)
    add_runs(
        s,
        x2,
        y,
        box_w,
        box_h,
        [
            ("Revenue-formel\n", 18, WHITE, True, FONT_HEAD),
            ("\n", 6, WHITE, False, FONT_BODY),
            ("N × $50 + Leads × $50\n", 16, GREEN, True, FONT_BODY),
            ("\n", 6, WHITE, False, FONT_BODY),
            ("$600/år subscription\n", 11, GRAY_LIGHT, False, FONT_BODY),
            ("+ $50 pr. quote-form-submit", 11, GRAY_LIGHT, False, FONT_BODY),
        ],
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
        line_spacing=1.15,
    )

    # Arrow 2
    a2 = s.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        x2 + box_w,
        arrow_y,
        gap,
        Inches(0.5),
    )
    fill(a2, GREEN)
    no_border(a2)

    # Box 3
    x3 = x2 + box_w + gap
    b3 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x3, y, box_w, box_h)
    fill(b3, WHITE)
    no_border(b3)
    add_runs(
        s,
        x3,
        y,
        box_w,
        box_h,
        [
            ("30.000\n", 56, GREEN_DARK, True, FONT_HEAD),
            ("DKK/måned i løn\n", 14, GRAY_TEXT, False, FONT_BODY),
            ("≈ $4.400/md fra partnere", 11, GRAY_TEXT, False, FONT_BODY),
        ],
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
        line_spacing=1.2,
    )

    # Footer caption
    add_text(
        s,
        Inches(0.5),
        Inches(5.8),
        Inches(12.333),
        Inches(0.6),
        "10k besøgende konverterer IKKE direkte til 40 partnere. Trafik driver leads ($50/stk) · outreach driver subscriptions ($600/år).",
        size=13,
        color=GRAY_TEXT,
        font=FONT_BODY,
        align=PP_ALIGN.CENTER,
    )


# =====================================================================
# Slide 3 — Regnestykket
# =====================================================================
def slide_math():
    s = blank_slide()
    add_header(s, "Regnestykket", 3)
    add_bg(s)

    # Headline
    add_text(
        s,
        Inches(0.5),
        Inches(1.55),
        Inches(12.333),
        Inches(0.5),
        "30.000 DKK/md  ≈  $4.400/md  (kurs 6,85 DKK/USD)",
        size=22,
        color=DARK_TEXT,
        bold=True,
        font=FONT_HEAD,
        align=PP_ALIGN.CENTER,
    )

    # Formula sub-headline
    add_text(
        s,
        Inches(0.5),
        Inches(2.1),
        Inches(12.333),
        Inches(0.4),
        "Formel:  Månedlig omsætning  =  N partnere × $50  +  Leads × $50",
        size=14,
        color=GREEN_DARK,
        bold=True,
        font=FONT_BODY,
        align=PP_ALIGN.CENTER,
    )

    # Three scenario cards
    card_w = Inches(3.7)
    card_h = Inches(3.6)
    gap = Inches(0.3)
    total_w = card_w * 3 + gap * 2
    start_x = (SLIDE_W - total_w) / 2
    y = Inches(2.9)

    scenarios = [
        {
            "label": "A · PARTNER-HEAVY",
            "big": "50",
            "small": "partnere + 38 leads/md",
            "calc": "50 × $50 + 38 × $50\n= $4.400/md",
            "note": "Kræver lift af conversion til 1,5%",
            "featured": False,
        },
        {
            "label": "B · MIX",
            "big": "40",
            "small": "partnere + 50 leads/md",
            "calc": "40 × $50 + 50 × $50\n= $4.500/md ≈ 30.800 DKK",
            "note": "Kræver 0,5% quote-conversion @ 10k MAU",
            "featured": True,
        },
        {
            "label": "C · LEAD-HEAVY",
            "big": "20",
            "small": "partnere + 70 leads/md",
            "calc": "20 × $50 + 70 × $50\n= $4.500/md",
            "note": "Få men aktive partnere, churn-risk",
            "featured": False,
        },
    ]

    for i, sc in enumerate(scenarios):
        x = start_x + (card_w + gap) * i
        if sc["featured"]:
            # Recommended badge
            badge = s.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x + Inches(0.8),
                y - Inches(0.25),
                card_w - Inches(1.6),
                Inches(0.4),
            )
            fill(badge, GREEN)
            no_border(badge)
            add_text(
                s,
                x + Inches(0.8),
                y - Inches(0.25),
                card_w - Inches(1.6),
                Inches(0.4),
                "ANBEFALET",
                size=11,
                color=WHITE,
                bold=True,
                font=FONT_BODY,
                align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE,
            )
            # Border via outer green rectangle
            outer = s.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h
            )
            fill(outer, GREEN)
            no_border(outer)
            inner = s.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x + Inches(0.07),
                y + Inches(0.07),
                card_w - Inches(0.14),
                card_h - Inches(0.14),
            )
            fill(inner, WHITE)
            no_border(inner)
        else:
            card = s.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h
            )
            fill(card, WHITE)
            no_border(card)

        # Label
        add_text(
            s,
            x,
            y + Inches(0.3),
            card_w,
            Inches(0.4),
            sc["label"],
            size=12,
            color=GREEN_DARK if sc["featured"] else GRAY_TEXT,
            bold=True,
            font=FONT_BODY,
            align=PP_ALIGN.CENTER,
        )

        # Big number
        add_text(
            s,
            x,
            y + Inches(0.8),
            card_w,
            Inches(1.2),
            sc["big"],
            size=64,
            color=GREEN_DARK,
            bold=True,
            font=FONT_HEAD,
            align=PP_ALIGN.CENTER,
        )

        # Small label
        add_text(
            s,
            x,
            y + Inches(2.0),
            card_w,
            Inches(0.4),
            sc["small"],
            size=13,
            color=DARK_TEXT,
            font=FONT_BODY,
            align=PP_ALIGN.CENTER,
        )

        # Calc
        add_text(
            s,
            x + Inches(0.2),
            y + Inches(2.5),
            card_w - Inches(0.4),
            Inches(0.7),
            sc["calc"],
            size=12,
            color=GRAY_TEXT,
            font=FONT_BODY,
            align=PP_ALIGN.CENTER,
        )

        # Note
        add_text(
            s,
            x + Inches(0.2),
            y + Inches(3.2),
            card_w - Inches(0.4),
            Inches(0.4),
            sc["note"],
            size=10,
            color=GRAY_TEXT,
            font=FONT_BODY,
            align=PP_ALIGN.CENTER,
        )

    # Footer
    add_text(
        s,
        Inches(0.5),
        Inches(6.7),
        Inches(12.333),
        Inches(0.4),
        "I dag: 1 partner × $600 banked + 0 leads. Gap til 30k DKK/md: 19-49 partnere + lead-conversion-lift.",
        size=12,
        color=RED_MUTED,
        bold=True,
        font=FONT_BODY,
        align=PP_ALIGN.CENTER,
    )


# =====================================================================
# Slide 4 — To-funnel-model
# =====================================================================
def slide_two_funnels():
    s = blank_slide()
    add_header(s, "To funnels — begge skal måles", 4)
    add_bg(s)

    # Two columns
    col_w = Inches(5.8)
    gap = Inches(0.5)
    start_x = (SLIDE_W - col_w * 2 - gap) / 2
    y = Inches(1.8)
    col_h = Inches(5.4)

    # Buyer funnel
    c1 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, start_x, y, col_w, col_h)
    fill(c1, WHITE)
    no_border(c1)
    add_text(
        s,
        start_x,
        y + Inches(0.2),
        col_w,
        Inches(0.5),
        "BUYER-FUNNEL",
        size=14,
        color=GREEN_DARK,
        bold=True,
        font=FONT_BODY,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s,
        start_x,
        y + Inches(0.7),
        col_w,
        Inches(0.5),
        "Trafik → Konvertering",
        size=20,
        color=DARK_TEXT,
        bold=True,
        font=FONT_HEAD,
        align=PP_ALIGN.CENTER,
    )

    buyer_steps = [
        ("Besøgende på platformen", "GA4 + sessions"),
        ("STL upload / quote start", "upload_events"),
        ("Supplier-views", "supplier_pageview"),
        ("Quote-form-submit  →  $50", "quote_request_submit = revenue"),
    ]
    step_y = y + Inches(1.4)
    for i, (label, src) in enumerate(buyer_steps):
        sy = step_y + Inches(0.85) * i
        box = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            start_x + Inches(0.4),
            sy,
            col_w - Inches(0.8),
            Inches(0.7),
        )
        fill(box, GRAY_LIGHT)
        no_border(box)
        add_runs(
            s,
            start_x + Inches(0.6),
            sy,
            col_w - Inches(1.0),
            Inches(0.7),
            [
                (f"{i+1}.  ", 14, GREEN_DARK, True, FONT_BODY),
                (label, 14, DARK_TEXT, True, FONT_BODY),
                (f"   ·   {src}", 11, GRAY_TEXT, False, FONT_BODY),
            ],
            anchor=MSO_ANCHOR.MIDDLE,
        )

    # Supplier funnel
    x2 = start_x + col_w + gap
    c2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x2, y, col_w, col_h)
    fill(c2, DARK)
    no_border(c2)
    add_text(
        s,
        x2,
        y + Inches(0.2),
        col_w,
        Inches(0.5),
        "SUPPLIER-FUNNEL",
        size=14,
        color=GREEN,
        bold=True,
        font=FONT_BODY,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s,
        x2,
        y + Inches(0.7),
        col_w,
        Inches(0.5),
        "Value → Revenue",
        size=20,
        color=WHITE,
        bold=True,
        font=FONT_HEAD,
        align=PP_ALIGN.CENTER,
    )

    supplier_steps = [
        ("Suppliers i directory", "suppliers count"),
        ("Leads leveret pr. supplier", "quote_submit + supplier_id (gap)"),
        ("Outreach + samtaler bookede", "manuelt (CRM)"),
        ("Underskriv $600/år-kontrakt", "is_partner=true · cash up front"),
    ]
    for i, (label, src) in enumerate(supplier_steps):
        sy = step_y + Inches(0.85) * i
        box = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x2 + Inches(0.4),
            sy,
            col_w - Inches(0.8),
            Inches(0.7),
        )
        fill(box, RGBColor(0x2A, 0x2A, 0x2A))
        no_border(box)
        add_runs(
            s,
            x2 + Inches(0.6),
            sy,
            col_w - Inches(1.0),
            Inches(0.7),
            [
                (f"{i+1}.  ", 14, GREEN, True, FONT_BODY),
                (label, 14, WHITE, True, FONT_BODY),
                (f"   ·   {src}", 11, GRAY_MED, False, FONT_BODY),
            ],
            anchor=MSO_ANCHOR.MIDDLE,
        )


# =====================================================================
# Slide 5 — Buyer KPI'er
# =====================================================================
def slide_buyer_kpis():
    s = blank_slide()
    add_header(s, "Buyer-KPI'er (ugentligt)", 5)
    add_bg(s)

    # Table
    headers = ["#", "KPI", "Kilde", "Ugemål @ 10k MAU"]
    col_widths = [Inches(0.6), Inches(5.0), Inches(4.0), Inches(2.7)]
    rows = [
        ["1", "Unikke besøgende", "GA4 / data-overview", "~2.500"],
        ["2", "STL-uploads", "upload_events", "~125  (5%)"],
        ["3", "Supplier-views", "supplier_pageview", "~500"],
        ["4", "Outbound clicks", "outbound_click", "~75"],
        ["5", "Quote-form-submits  ⇒  $50/stk", "quote_request_submit", "~12  ($600/uge)"],
        ["6", "Quote-form conversion-rate", "beregnet (submits/visits)", "0,5%  (i dag ~0,25%)"],
    ]

    table_x = Inches(0.5)
    table_y = Inches(1.9)
    header_h = Inches(0.55)
    row_h = Inches(0.6)

    # Header row
    cx = table_x
    for w, h in zip(col_widths, headers):
        cell = add_rect(s, cx, table_y, w, header_h, DARK)
        add_text(
            s,
            cx,
            table_y,
            w,
            header_h,
            h,
            size=12,
            color=WHITE,
            bold=True,
            font=FONT_BODY,
            align=PP_ALIGN.LEFT if h != "#" else PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        cx += w

    # Data rows
    for i, row in enumerate(rows):
        ry = table_y + header_h + row_h * i
        bg_color = WHITE if i % 2 == 0 else GRAY_LIGHT
        # row background
        rx = table_x
        for w in col_widths:
            add_rect(s, rx, ry, w, row_h, bg_color)
            rx += w

        # Highlight rows 5 and 6 (quote-submit + conv-rate — revenue-critical KPIs)
        cx = table_x
        for j, (w, val) in enumerate(zip(col_widths, row)):
            is_last = j == len(row) - 1
            is_headline = i in (4, 5)
            text_color = GREEN_DARK if is_last else DARK_TEXT
            if is_headline:
                text_color = GREEN_DARK
            bold = is_last or is_headline
            align = PP_ALIGN.LEFT
            if j == 0:
                align = PP_ALIGN.CENTER
            add_text(
                s,
                cx + Inches(0.1),
                ry,
                w - Inches(0.2),
                row_h,
                val,
                size=13 if not is_headline else 14,
                color=text_color,
                bold=bold,
                font=FONT_BODY,
                align=align,
                anchor=MSO_ANCHOR.MIDDLE,
            )
            cx += w

    # Footer note
    add_text(
        s,
        Inches(0.5),
        Inches(6.6),
        Inches(12.333),
        Inches(0.5),
        "KPI #5 = revenue-event ($50 pr. submit). KPI #6 = højeste leverage: lift 0,25% → 0,5% = 2× lead-revenue uden ny trafik.",
        size=12,
        color=GREEN_DARK,
        bold=True,
        font=FONT_BODY,
        align=PP_ALIGN.LEFT,
    )


# =====================================================================
# Slide 6 — Supplier KPI'er
# =====================================================================
def slide_supplier_kpis():
    s = blank_slide()
    add_header(s, "Supplier-KPI'er (revenue path)", 6)
    add_bg(s)

    headers = ["#", "KPI", "Kilde i dag", "Status"]
    col_widths = [Inches(0.6), Inches(5.5), Inches(4.0), Inches(2.2)]
    rows = [
        ["7", "Aktive suppliers i directory", "suppliers count", "TRACKES"],
        ["8", "Leads leveret pr. partner (uge)", "quote_submit + supplier_id", "GAP"],
        ["9", "Booked partner-samtaler (uge)", "manuelt (CRM/regneark)", "GAP"],
        ["10", "Nye partnere (uge) × $600", "is_partner=true", "TRACKES"],
        ["11", "Annual revenue booked (kvt.)", "$600 × nye partnere", "GAP"],
        ["12", "Renewal-rate (efter 12 mdr.)", "fra md. 13", "FREMTID"],
    ]

    table_x = Inches(0.5)
    table_y = Inches(1.9)
    header_h = Inches(0.55)
    row_h = Inches(0.6)

    cx = table_x
    for w, h in zip(col_widths, headers):
        add_rect(s, cx, table_y, w, header_h, DARK)
        add_text(
            s,
            cx,
            table_y,
            w,
            header_h,
            h,
            size=12,
            color=WHITE,
            bold=True,
            font=FONT_BODY,
            align=PP_ALIGN.LEFT if h != "#" else PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        cx += w

    for i, row in enumerate(rows):
        ry = table_y + header_h + row_h * i
        bg_color = WHITE if i % 2 == 0 else GRAY_LIGHT
        rx = table_x
        for w in col_widths:
            add_rect(s, rx, ry, w, row_h, bg_color)
            rx += w

        cx = table_x
        for j, (w, val) in enumerate(zip(col_widths, row)):
            is_status = j == 3
            is_revenue = row[0] in ("10", "11")
            text_color = DARK_TEXT
            bold = False
            if is_status:
                if val == "GAP":
                    text_color = RED_MUTED
                    bold = True
                elif val == "FREMTID":
                    text_color = AMBER
                    bold = True
                else:
                    text_color = GREEN_DARK
                    bold = True
            elif is_revenue and j == 1:
                text_color = GREEN_DARK
                bold = True

            align = PP_ALIGN.LEFT
            if j == 0:
                align = PP_ALIGN.CENTER
            if is_status:
                align = PP_ALIGN.CENTER

            add_text(
                s,
                cx + Inches(0.1),
                ry,
                w - Inches(0.2),
                row_h,
                val,
                size=13,
                color=text_color,
                bold=bold,
                font=FONT_BODY,
                align=align,
                anchor=MSO_ANCHOR.MIDDLE,
            )
            cx += w

    # Footer
    add_text(
        s,
        Inches(0.5),
        Inches(6.6),
        Inches(12.333),
        Inches(0.5),
        "3 GAPs + 1 fremtid. Lead-attribution (KPI #8) er revenue-kritisk — hver utracket submit = $50 tabt. Se slide 10.",
        size=12,
        color=RED_MUTED,
        bold=True,
        font=FONT_BODY,
    )


# =====================================================================
# Slide 7 — Baseline i dag
# =====================================================================
def slide_baseline():
    s = blank_slide()
    add_header(s, "Baseline i dag", 7)
    add_bg(s)

    add_text(
        s,
        Inches(0.5),
        Inches(1.65),
        Inches(12.333),
        Inches(0.5),
        "Udfyldes hver mandag fra /admin/data-overview (7-dages vindue)",
        size=14,
        color=GRAY_TEXT,
        font=FONT_BODY,
    )

    # Two columns of metric cards
    metrics = [
        ("Unikke besøgende/uge", "____", "mål: 2.500", False),
        ("Quote-form-submits/uge", "____", "mål: 12 ($600/uge)", False),
        ("Conversion-rate", "____", "mål: 0,5%", False),
        ("Outbound clicks/uge", "____", "mål: 75", False),
        ("Aktive suppliers", "____", "directory", False),
        ("Betalende partnere", "1", "AM Printservice", True),
        ("Revenue booked YTD", "$600", "1 × $600/år", True),
        ("Gap til 30k DKK/md", "$4.350", "≈ 29.800 DKK", True),
    ]

    card_w = Inches(3.0)
    card_h = Inches(1.7)
    gap_x = Inches(0.2)
    gap_y = Inches(0.2)
    cols = 4
    start_x = (SLIDE_W - card_w * cols - gap_x * (cols - 1)) / 2
    y = Inches(2.3)

    for i, (label, val, sub, known) in enumerate(metrics):
        row_idx = i // cols
        col_idx = i % cols
        x = start_x + col_idx * (card_w + gap_x)
        cy = y + row_idx * (card_h + gap_y)

        bg_color = WHITE if known else GRAY_LIGHT
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, cy, card_w, card_h)
        fill(card, bg_color)
        no_border(card)

        # Side stripe
        stripe_color = GREEN if known else AMBER
        stripe = s.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, cy, Inches(0.12), card_h
        )
        fill(stripe, stripe_color)
        no_border(stripe)

        add_text(
            s,
            x + Inches(0.3),
            cy + Inches(0.15),
            card_w - Inches(0.4),
            Inches(0.35),
            label,
            size=11,
            color=GRAY_TEXT,
            bold=True,
            font=FONT_BODY,
        )
        add_text(
            s,
            x + Inches(0.3),
            cy + Inches(0.5),
            card_w - Inches(0.4),
            Inches(0.8),
            val,
            size=36,
            color=GREEN_DARK if known else GRAY_TEXT,
            bold=True,
            font=FONT_HEAD,
        )
        add_text(
            s,
            x + Inches(0.3),
            cy + Inches(1.3),
            card_w - Inches(0.4),
            Inches(0.3),
            sub,
            size=10,
            color=GRAY_TEXT,
            font=FONT_BODY,
        )

    # Legend
    add_runs(
        s,
        Inches(0.5),
        Inches(6.8),
        Inches(12.333),
        Inches(0.4),
        [
            ("■ ", 14, GREEN, True, FONT_BODY),
            ("kendt baseline   ", 11, GRAY_TEXT, False, FONT_BODY),
            ("■ ", 14, AMBER, True, FONT_BODY),
            ("udfyldes ugentligt fra admin", 11, GRAY_TEXT, False, FONT_BODY),
        ],
    )


# =====================================================================
# Slide 8 — Conversion-rate hierarki
# =====================================================================
def slide_funnel_pyramid():
    s = blank_slide()
    add_header(s, "Konvertering — fra trafik til kliks", 8)
    add_bg(s)

    # Stacked horizontal funnel bars
    stages = [
        ("Besøgende", 2500, 100, GREEN_DARK),
        ("STL-upload eller quote-start", 250, 10, GREEN),
        ("Supplier-views", 500, 20, GREEN),  # note: many views per session
        ("Outbound clicks", 75, 3, AMBER),
        ("Quote requests", 25, 1, RED_MUTED),
    ]

    # Recalculate to make visual sense (pyramid shape)
    # Use absolute counts as widths
    max_count = 2500
    bar_max_w = Inches(9.0)
    bar_h = Inches(0.7)
    gap = Inches(0.25)
    start_y = Inches(2.0)
    start_x = Inches(2.2)

    for i, (label, count, pct, color) in enumerate(stages):
        bar_w = bar_max_w * (count / max_count)
        if bar_w < Inches(1.5):
            bar_w = Inches(1.5)  # min width for legibility
        by = start_y + (bar_h + gap) * i
        # Center the bar
        bx = start_x + (bar_max_w - bar_w) / 2

        bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, by, bar_w, bar_h)
        fill(bar, color)
        no_border(bar)

        # Label inside bar
        add_text(
            s,
            bx,
            by,
            bar_w,
            bar_h,
            f"{count:,}".replace(",", "."),
            size=18,
            color=WHITE,
            bold=True,
            font=FONT_HEAD,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )

        # Stage label to the left
        add_text(
            s,
            Inches(0.4),
            by,
            Inches(1.7),
            bar_h,
            label,
            size=13,
            color=DARK_TEXT,
            bold=True,
            font=FONT_BODY,
            align=PP_ALIGN.RIGHT,
            anchor=MSO_ANCHOR.MIDDLE,
        )

        # Pct to the right
        add_text(
            s,
            Inches(11.3),
            by,
            Inches(1.6),
            bar_h,
            f"{pct}%",
            size=14,
            color=GREEN_DARK,
            bold=True,
            font=FONT_BODY,
            align=PP_ALIGN.LEFT,
            anchor=MSO_ANCHOR.MIDDLE,
        )

    # Note about supplier-views
    add_text(
        s,
        Inches(0.5),
        Inches(6.4),
        Inches(12.333),
        Inches(0.5),
        "* Supplier-views = 20% pga. flere visninger pr. session (typisk 4-5 leverandører i én sammenligning).",
        size=11,
        color=GRAY_TEXT,
        font=FONT_BODY,
    )
    add_text(
        s,
        Inches(0.5),
        Inches(6.85),
        Inches(12.333),
        Inches(0.5),
        "Headline-revenue: hver Quote-form-submit = $50. Lift 0,25% → 0,5% = $1.250/md ekstra uden ny trafik.",
        size=12,
        color=GREEN_DARK,
        bold=True,
        font=FONT_BODY,
    )


# =====================================================================
# Slide 9 — Ugentlig review-rutine
# =====================================================================
def slide_routine():
    s = blank_slide()
    add_header(s, "Ugentlig review-rutine", 9)
    add_bg(s)

    # When/where banner
    banner = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.7),
        Inches(1.8),
        Inches(11.9),
        Inches(0.9),
    )
    fill(banner, DARK)
    no_border(banner)
    add_runs(
        s,
        Inches(1.0),
        Inches(1.8),
        Inches(11.3),
        Inches(0.9),
        [
            ("MANDAG MORGEN   ·   30 MIN   ·   ", 16, GREEN, True, FONT_BODY),
            ("/admin/data-overview  +  dette KPI-deck", 16, WHITE, False, FONT_BODY),
        ],
        anchor=MSO_ANCHOR.MIDDLE,
    )

    # 4 steps
    steps = [
        ("1", "Notér tallene", "For hver KPI: denne uge vs. sidste uge vs. 4-ugers rullende gennemsnit."),
        ("2", "Find bevægelsen", "Hvilke 1-2 metrics rykker mest? Skriv ÉN linje om hvorfor."),
        ("3", "Vælg ÉT fokus", "Top-of-funnel, conversion, eller supplier-outreach? Skriv det ned."),
        ("4", "Log beslutningen", "Tilføj en ugeslide eller en linje i Notion. Hold dig selv accountable."),
    ]
    step_w = Inches(2.85)
    step_h = Inches(3.4)
    gap = Inches(0.2)
    start_x = (SLIDE_W - step_w * 4 - gap * 3) / 2
    y = Inches(3.0)

    for i, (num, head, body) in enumerate(steps):
        x = start_x + i * (step_w + gap)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, step_w, step_h)
        fill(card, WHITE)
        no_border(card)

        # Big number circle
        circ = s.shapes.add_shape(
            MSO_SHAPE.OVAL,
            x + (step_w - Inches(1)) / 2,
            y + Inches(0.4),
            Inches(1),
            Inches(1),
        )
        fill(circ, GREEN)
        no_border(circ)
        add_text(
            s,
            x + (step_w - Inches(1)) / 2,
            y + Inches(0.4),
            Inches(1),
            Inches(1),
            num,
            size=38,
            color=WHITE,
            bold=True,
            font=FONT_HEAD,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )

        add_text(
            s,
            x + Inches(0.2),
            y + Inches(1.7),
            step_w - Inches(0.4),
            Inches(0.5),
            head,
            size=18,
            color=DARK_TEXT,
            bold=True,
            font=FONT_HEAD,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            s,
            x + Inches(0.3),
            y + Inches(2.3),
            step_w - Inches(0.6),
            Inches(1.0),
            body,
            size=12,
            color=GRAY_TEXT,
            font=FONT_BODY,
            align=PP_ALIGN.CENTER,
        )


# =====================================================================
# Slide 10 — Instrumentation gaps
# =====================================================================
def slide_gaps():
    s = blank_slide()
    add_header(s, "Hvad mangler vi at tracke?", 10)
    add_bg(s)

    add_text(
        s,
        Inches(0.5),
        Inches(1.65),
        Inches(12.333),
        Inches(0.5),
        "3 instrumentations-bygges for at KPI'erne kan måles løbende.",
        size=14,
        color=GRAY_TEXT,
        font=FONT_BODY,
    )

    gaps = [
        {
            "title": "partner_subscriptions tabel",
            "drives": "Driver KPI #10/#11 (Revenue booked + Renewal)",
            "details": "Simpel: supplier_id, paid_amount_usd (600), paid_at, expires_at (+12md), renewal_status. Ingen tiers. Manuel fakturering for nu, Stripe senere.",
        },
        {
            "title": "Lead-attribution (revenue-kritisk)",
            "drives": "Driver KPI #5 ($50/lead) + KPI #8",
            "details": "Tilføj supplier_id til quote_request_submit-event + ny partner_leads tabel. Hver utracket submit = $50 tabt. Samtidig fix: selectItems/searches/supplier_id GTM-tagging.",
        },
        {
            "title": "MAU-query i admin",
            "drives": "Driver KPI #1 (Unikke besøgende)",
            "details": "Brug GA4-edge function (ga4-analytics) eller udled fra analytics_events.session_id distinct count over 30 dage. Vis i /admin/data-overview.",
        },
    ]

    card_w = Inches(4.0)
    card_h = Inches(4.5)
    gap = Inches(0.2)
    start_x = (SLIDE_W - card_w * 3 - gap * 2) / 2
    y = Inches(2.3)

    for i, gap_data in enumerate(gaps):
        x = start_x + i * (card_w + gap)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h)
        fill(card, WHITE)
        no_border(card)

        # Red top stripe
        stripe = s.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, y, card_w, Inches(0.15)
        )
        fill(stripe, RED_MUTED)
        no_border(stripe)

        # GAP label
        gap_label = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x + Inches(0.3),
            y + Inches(0.4),
            Inches(0.9),
            Inches(0.35),
        )
        fill(gap_label, RED_MUTED)
        no_border(gap_label)
        add_text(
            s,
            x + Inches(0.3),
            y + Inches(0.4),
            Inches(0.9),
            Inches(0.35),
            "GAP",
            size=11,
            color=WHITE,
            bold=True,
            font=FONT_BODY,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )

        # Title
        add_text(
            s,
            x + Inches(0.3),
            y + Inches(0.9),
            card_w - Inches(0.6),
            Inches(0.9),
            gap_data["title"],
            size=20,
            color=DARK_TEXT,
            bold=True,
            font=FONT_HEAD,
        )

        # Drives
        add_text(
            s,
            x + Inches(0.3),
            y + Inches(1.9),
            card_w - Inches(0.6),
            Inches(0.5),
            gap_data["drives"],
            size=12,
            color=GREEN_DARK,
            bold=True,
            font=FONT_BODY,
        )

        # Details
        add_text(
            s,
            x + Inches(0.3),
            y + Inches(2.5),
            card_w - Inches(0.6),
            Inches(1.8),
            gap_data["details"],
            size=12,
            color=GRAY_TEXT,
            font=FONT_BODY,
        )


# =====================================================================
# Slide 11 — 90-dages plan
# =====================================================================
def slide_90_day_plan():
    s = blank_slide()
    add_header(s, "90-dages fokus", 11)
    add_bg(s)

    add_text(
        s,
        Inches(0.5),
        Inches(1.65),
        Inches(12.333),
        Inches(0.5),
        "Næste 12 uger: ÉN KPI er det primære fokus pr. uge. Skift uge for uge.",
        size=14,
        color=GRAY_TEXT,
        font=FONT_BODY,
    )

    weeks = [
        ("Uge 1-3", "LEAD-ATTRIBUTION", "Tilføj supplier_id til quote_submit · partner_leads tabel · partner_subscriptions tabel · MAU-query", DARK, WHITE),
        ("Uge 4-6", "BUYER CONVERSION", "Quote-form-UX · CTA-tests · STL-flow · lift quote-submit fra 0,25% → 0,5% = 2× lead-revenue", GREEN_DARK, WHITE),
        ("Uge 7-9", "BUYER TOP-OF-FUNNEL", "SEO-content · Compare-prices-side · landing pages · hæv MAU fra X → 10k for mere lead-volumen", GREEN, WHITE),
        ("Uge 10-12", "SUPPLIER OUTREACH", "Outreach til top-25 suppliers · vis dem leads-rapport · få 3 nye $600/år-kontrakter underskrevet", AMBER, WHITE),
    ]

    block_w = Inches(2.95)
    block_h = Inches(4.4)
    gap = Inches(0.15)
    start_x = (SLIDE_W - block_w * 4 - gap * 3) / 2
    y = Inches(2.3)

    for i, (period, focus, detail, bg, fg) in enumerate(weeks):
        x = start_x + i * (block_w + gap)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, block_w, block_h)
        fill(card, bg)
        no_border(card)

        add_text(
            s,
            x + Inches(0.3),
            y + Inches(0.4),
            block_w - Inches(0.6),
            Inches(0.5),
            period,
            size=14,
            color=fg,
            bold=True,
            font=FONT_BODY,
            align=PP_ALIGN.CENTER,
        )

        # Divider
        div = s.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x + Inches(0.6),
            y + Inches(1.0),
            block_w - Inches(1.2),
            Inches(0.04),
        )
        fill(div, fg)
        no_border(div)

        add_text(
            s,
            x + Inches(0.2),
            y + Inches(1.3),
            block_w - Inches(0.4),
            Inches(1.0),
            focus,
            size=20,
            color=fg,
            bold=True,
            font=FONT_HEAD,
            align=PP_ALIGN.CENTER,
        )

        add_text(
            s,
            x + Inches(0.3),
            y + Inches(2.7),
            block_w - Inches(0.6),
            Inches(1.5),
            detail,
            size=11,
            color=fg,
            font=FONT_BODY,
            align=PP_ALIGN.CENTER,
        )

    # Target outcome
    target = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.7),
        Inches(6.85),
        Inches(11.9),
        Inches(0.5),
    )
    fill(target, GREEN_DARK)
    no_border(target)
    add_text(
        s,
        Inches(1.0),
        Inches(6.85),
        Inches(11.3),
        Inches(0.5),
        "EFTER 90 DAGE:  4 nye partnere × $600 = $2.400 booked  ·  0,5% quote-conversion  ·  ~$1.250/md lead-revenue",
        size=12,
        color=WHITE,
        bold=True,
        font=FONT_BODY,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )


# =====================================================================
# Slide 12 — Datakilder (bagsidekort)
# =====================================================================
def slide_sources():
    s = blank_slide()
    add_header(s, "Datakilder", 12)
    add_bg(s)

    add_text(
        s,
        Inches(0.5),
        Inches(1.65),
        Inches(12.333),
        Inches(0.5),
        "Hvor henter jeg hver KPI fra — kode-referencer + dashboards.",
        size=14,
        color=GRAY_TEXT,
        font=FONT_BODY,
    )

    sources = [
        ("KPI #1-6 (Buyer-funnel)", [
            "src/lib/analytics.ts:13-30 — HIGH_SIGNAL_EVENTS",
            "src/hooks/use-funnel-data.ts:29-83 — funnel-beregning",
            "src/hooks/use-ga4-funnel.ts — GA4 integration",
            "Dashboard: /admin/data-overview",
        ]),
        ("KPI #7-10 (Supplier-funnel)", [
            "suppliers tabel — is_partner boolean kolonne",
            "src/pages/admin/Admin.tsx — supplier inventory panel",
            "Outreach-CRM: manuelt regneark (TBD)",
            "Dashboard: /admin (Suppliers-fane)",
        ]),
        ("KPI #5/#10/#11 (Revenue) — BYG", [
            "partner_subscriptions: supplier_id, paid_at, expires_at, $600",
            "partner_leads: quote_event_id → supplier_id → $50",
            "Manuel fakturering (Stripe senere)",
            "Vises på: /admin/data-overview (ny sektion)",
        ]),
    ]

    card_w = Inches(4.0)
    card_h = Inches(4.5)
    gap = Inches(0.2)
    start_x = (SLIDE_W - card_w * 3 - gap * 2) / 2
    y = Inches(2.3)

    for i, (title, items) in enumerate(sources):
        x = start_x + i * (card_w + gap)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h)
        fill(card, WHITE)
        no_border(card)

        # Top color stripe
        stripe_color = [GREEN, GREEN_DARK, AMBER][i]
        stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, card_w, Inches(0.15))
        fill(stripe, stripe_color)
        no_border(stripe)

        add_text(
            s,
            x + Inches(0.3),
            y + Inches(0.4),
            card_w - Inches(0.6),
            Inches(0.6),
            title,
            size=16,
            color=DARK_TEXT,
            bold=True,
            font=FONT_HEAD,
        )

        for j, item in enumerate(items):
            iy = y + Inches(1.2) + Inches(0.7) * j
            # bullet
            bullet = s.shapes.add_shape(
                MSO_SHAPE.OVAL,
                x + Inches(0.4),
                iy + Inches(0.15),
                Inches(0.12),
                Inches(0.12),
            )
            fill(bullet, stripe_color)
            no_border(bullet)
            add_text(
                s,
                x + Inches(0.65),
                iy,
                card_w - Inches(0.85),
                Inches(0.6),
                item,
                size=11,
                color=GRAY_TEXT,
                font=FONT_BODY,
            )


# =====================================================================
# Build deck
# =====================================================================
def main():
    slide_title()
    slide_north_star()
    slide_math()
    slide_two_funnels()
    slide_buyer_kpis()
    slide_supplier_kpis()
    slide_baseline()
    slide_funnel_pyramid()
    slide_routine()
    slide_gaps()
    slide_90_day_plan()
    slide_sources()
    prs.save(str(OUTPUT))
    print(f"Wrote: {OUTPUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
